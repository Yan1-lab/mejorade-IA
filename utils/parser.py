# utils/parser.py
import os
from pathlib import Path
from pypdf import PdfReader
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import pandas as pd

UPLOADS = Path("uploads")
UPLOADS.mkdir(exist_ok=True)

def save_uploaded_file(uploaded_file, user_id:int, session_id:int):
    user_dir = UPLOADS / str(user_id)
    session_dir = user_dir / str(session_id)
    session_dir.mkdir(parents=True, exist_ok=True)
    filepath = session_dir / uploaded_file.name
    with open(filepath, "wb") as f:
        f.write(uploaded_file.getbuffer())
    ext = uploaded_file.name.split(".")[-1].lower()
    extracted = extract_text(str(filepath), ext)
    return {"filename": uploaded_file.name, "filepath": str(filepath), "filetype": ext, "extracted_text": extracted}

def extract_text(path: str, ext: str) -> str:
    ext = ext.lower()
    try:
        if ext == "pdf":
            return extract_text_pdf(path)
        if ext in ("png","jpg","jpeg","bmp","tiff","gif"):
            return extract_text_image(path)
        if ext in ("docx","doc"):
            return extract_text_docx(path)
        if ext in ("xlsx","xls","csv"):
            return extract_text_excel(path)
        if ext in ("pptx","ppt"):
            return extract_text_pptx(path)
        if ext in ("txt","py","js","java","c","cpp","rb","go","rs"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        return ""
    except Exception as e:
        return f"[EXTRACT ERROR: {e}]"

def extract_text_pdf(path):
    try:
        reader = PdfReader(path)
        text = ""
        for p in reader.pages:
            t = p.extract_text()
            if t:
                text += t + "\n"
        if text.strip():
            return text
        # fallback OCR
        images = convert_from_bytes(open(path,"rb").read())
        out = ""
        for im in images:
            out += pytesseract.image_to_string(im)
        return out
    except Exception as e:
        return f"[PDF ERROR: {e}]"

def extract_text_image(path):
    try:
        im = Image.open(path)
        return pytesseract.image_to_string(im)
    except Exception as e:
        return f"[IMG OCR ERROR: {e}]"

def extract_text_docx(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"[DOCX ERROR: {e}]"

def extract_text_excel(path):
    try:
        df = pd.read_excel(path, engine='openpyxl')
        return df.to_csv(index=False)
    except Exception as e:
        return f"[EXCEL ERROR: {e}]"

def extract_text_pptx(path):
    try:
        prs = Presentation(path)
        texts=[]
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[PPTX ERROR: {e}]"
